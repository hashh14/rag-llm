import os
import time
import streamlit as st
from langchain_community.vectorstores import Chroma
from langchain.prompts import ChatPromptTemplate
from langchain_community.llms.ollama import Ollama
from pptx import Presentation
from pptx.util import Pt
from get_embedding_function import get_embedding_function
from populate_database import load_documents, split_documents, add_to_chroma, delete_database

# Constants
CHROMA_PATH = "chroma"
DATA_PATH = "data"
UPLOAD_PATH = "data"
PPTX_FILE_PATH = "output_presentation.pptx"
MAX_CHARS_PER_SLIDE = 600

# Prompt templates
PROMPT_TEMPLATES = {
    "rag": """
    Answer the question based only on the following context:
    {context}
    ---
    Answer the question based on the above context: {question}
    """,
    "sum": """
    Summarize the given topic based on the provided context. 
    The summary must be concise, clear, and must cover all points specified by the user. 
    If a specific length is mentioned by the user, ensure the summary adheres to that length. 
    If no length is specified, create a summary that is appropriate in length for the given topic and provided context, without missing any important details.
    Context:
    {context}
    ---
    Topic:
    {question}
    """,
    "qa": """
    You are a question generating and answering model.
    Generate questions for the user based on the given query and context, and generate appropriate answers for all the questions that you generate.
    Generate questions and answers based on the type of question that the user specifies.
    The different types of questions can be:
    - Multiple choice questions or MCQ (Default should be 4 choices for each question, if not specified by the user. Mention the correct option by printing (Correct option:) at the end. By default, give 5 MCQs if not explicitly specified by the user.),
    - Fill in the blanks (where a blank is left in the sentence and the correct answer is present inside parentheses () at the end of the sentence. By default, give 5 fill in the blank questions if not explicitly specified by the user.),
    - True or False questions (where the question statement is given and the correct answer which is true or false, is given inside parentheses () at the end of the sentence. By default, give 5 true or false questions if not explicitly specified by the user.),
    - Or it can be any other type of question that the user specifies.
    Context:
    {context}
    ---
    Query:
    {question}
    """,
    "case": """
    You are a case study generation model.
    Based on the given query and the corresponding provided context, create a case study for the user.
    The case study must take some example related to the topic and context and explain the concept to the user using that example. 
    The example can be based on real events or can be a made-up story. 
    The case study must also have some conclusion at the end which explains the need for the solution or the concept based on the given context.
    Context:
    {context}
    ---
    Query:
    {question}
    """
}

def setup_sidebar():
    st.sidebar.title("Database Management")
    uploaded_files = st.sidebar.file_uploader("Upload Files", accept_multiple_files=True)
    
    if uploaded_files:
        for uploaded_file in uploaded_files:
            file_path = os.path.join(UPLOAD_PATH, uploaded_file.name)
            with open(file_path, "wb") as f:
                f.write(uploaded_file.getbuffer())
            success_message = st.sidebar.success(f"Saved file: {uploaded_file.name}")
            time.sleep(1)
            success_message.empty()

    colb1, colb2 = st.sidebar.columns(2)
    if colb1.button("Update Vector Database"):
        update_vector_database()
    if colb2.button("Reset Vector Database"):
        reset_vector_database()

    display_uploaded_files()
    setup_download_section()

def update_vector_database():
    documents = load_documents()
    chunks = split_documents(documents)
    add_to_chroma(chunks)
    success_message = st.sidebar.success("Vector database updated successfully.")
    time.sleep(2)
    success_message.empty()

def reset_vector_database():
    documents = load_documents()
    chunks = split_documents(documents)
    delete_database(chunks)
    documents = load_documents()
    chunks = split_documents(documents)
    add_to_chroma(chunks)
    success_message = st.sidebar.success("Vector database reset successfully.")
    time.sleep(2)
    success_message.empty()

def display_uploaded_files():
    st.sidebar.title("Uploaded Files")
    files = os.listdir(UPLOAD_PATH)
    if files:
        for file in files:
            file_path = os.path.join(UPLOAD_PATH, file)
            col1, col2 = st.sidebar.columns([4, 1])
            col1.write(file)
            if col2.button("Delete", key=file):
                os.remove(file_path)
                st.sidebar.warning(f"Deleted file: {file}")
                time.sleep(1)
                st.rerun()
    else:
        st.sidebar.write("No files available")

def setup_download_section():
    st.sidebar.title("Download Files")
    if not os.path.exists(PPTX_FILE_PATH):
        st.sidebar.write("Start chatting to download the PowerPoint file.")
    else:
        with open(PPTX_FILE_PATH, "rb") as f:
            pptx_bytes = f.read()
            st.sidebar.download_button(
                label="PowerPoint File",
                data=pptx_bytes,
                file_name=PPTX_FILE_PATH,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

def query_database(query_text: str, prompt_type: str):
    embedding_function = get_embedding_function()
    db = Chroma(persist_directory=CHROMA_PATH, embedding_function=embedding_function)
    results = db.similarity_search_with_score(query_text, k=5)
    
    context_text = "\n\n---\n\n".join([doc.page_content for doc, _score in results])
    prompt_template = ChatPromptTemplate.from_template(PROMPT_TEMPLATES[prompt_type])
    prompt = prompt_template.format(context=context_text, question=query_text)
    
    model = Ollama(model="mistral")
    response_text = model.invoke(prompt)
    
    sources = [doc.metadata.get("id", None) for doc, _score in results]
    update_presentation(query_text, response_text)
    
    return response_text

def update_presentation(title, content):
    if os.path.exists(PPTX_FILE_PATH):
        prs = Presentation(PPTX_FILE_PATH)
    else:
        prs = Presentation()
    
    add_slide_with_content(prs, title, content)
    prs.save(PPTX_FILE_PATH)

def add_slide_with_content(prs, title, content):
    slide_layout = prs.slide_layouts[1]
    special_cases = ["summary", "question", "questions", "mcq", "fill in the blanks", "fill in the blank", "case study", "case studies"]
    
    if any(case in title.lower() for case in special_cases):
        if "summary" in title.lower():
            title = "Summary:"
        elif "case study" in title.lower() or "case studies" in title.lower():
            title = "Case study:"
        else:
            title = "Q&A:"
    
    paragraphs = content.split("\n")
    current_content = ""
    first_slide = True
    
    for paragraph in paragraphs:
        if len(current_content) + len(paragraph) > MAX_CHARS_PER_SLIDE:
            add_content_to_slide(prs, slide_layout, title, current_content, first_slide)
            first_slide = False
            current_content = paragraph
        else:
            if current_content:
                current_content += "\n"
            current_content += paragraph
    
    add_content_to_slide(prs, slide_layout, title, current_content, first_slide)

def add_content_to_slide(prs, slide_layout, title, content, is_first_slide):
    slide = prs.slides.add_slide(slide_layout)
    content_placeholder = slide.placeholders[1]
    
    if is_first_slide:
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
    
    text_frame = content_placeholder.text_frame
    text_frame.word_wrap = True
    p = text_frame.add_paragraph()
    p.text = content
    p.font.size = Pt(20)

def main():
    st.title("Document Query Engine")
    setup_sidebar()
    
    function = st.selectbox("Select Function", ["Chat", "Summary Generation", "Question Generation", "Case Study Generation"])
    
    if 'messages' not in st.session_state:
        st.session_state.messages = []
    
    query_text = st.chat_input("Enter Query Text:")
    
    if query_text:
        function_map = {
            "Chat": "rag",
            "Summary Generation": "sum",
            "Question Generation": "qa",
            "Case Study Generation": "case"
        }
        response_text = query_database(query_text, function_map[function])
        st.session_state.messages.append({"role": "user", "content": query_text})
        st.session_state.messages.append({"role": "bot", "content": response_text})
    
    for message in st.session_state.messages:
        st.write(f"{message['role'].capitalize()}: {message['content']}")

if __name__ == "__main__":
    main()